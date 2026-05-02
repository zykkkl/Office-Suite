from pathlib import Path

PROJECT_ROOT = Path(__file__).parent.parent

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from office_suite.dsl.parser import parse_yaml_string
from office_suite.ir.compiler import compile_document
from office_suite.renderer.pptx.deck import PPTXRenderer


def test_pptx_text_layout_fields():
    dsl = """
version: "4.0"
type: presentation
theme: default
slides:
  - layout: blank
    elements:
      - type: text
        content: "1"
        position: { x: 30mm, y: 30mm, width: 16mm, height: 16mm }
        text_align: center
        verticalAlign: middle
        margin: 0
        style:
          font: { family: "Arial", size: 15, weight: 700, color: "#FFFFFF" }
"""
    out = Path(__file__).parent / "output" / "text_layout.pptx"
    doc = parse_yaml_string(dsl)
    ir = compile_document(doc)

    PPTXRenderer().render(ir, out)

    prs = Presentation(out)
    shape = next(s for s in prs.slides[0].shapes if s.has_text_frame and s.text.strip() == "1")
    text_frame = shape.text_frame
    paragraph = text_frame.paragraphs[0]

    assert paragraph.alignment == PP_ALIGN.CENTER
    assert text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE
    assert text_frame.margin_left == 0
    assert text_frame.margin_right == 0
    assert text_frame.margin_top == 0
    assert text_frame.margin_bottom == 0


def test_pptx_text_wrap_false_field():
    dsl = """
version: "4.0"
type: presentation
theme: default
slides:
  - layout: blank
    elements:
      - type: text
        content: "02 / 10"
        position: { x: 196mm, y: 126mm, width: 38mm, height: 8mm }
        align: right
        vertical_align: middle
        margin: 0
        wrap: false
        style:
          font: { family: "Arial", size: 13, color: "#1D2939" }
"""
    out = Path(__file__).parent / "output" / "text_wrap_false.pptx"
    doc = parse_yaml_string(dsl)
    ir = compile_document(doc)

    PPTXRenderer().render(ir, out)

    prs = Presentation(out)
    shape = next(s for s in prs.slides[0].shapes if s.has_text_frame and s.text.strip() == "02 / 10")
    text_frame = shape.text_frame
    paragraph = text_frame.paragraphs[0]

    assert paragraph.alignment == PP_ALIGN.RIGHT
    assert text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE
    assert text_frame.word_wrap is False
